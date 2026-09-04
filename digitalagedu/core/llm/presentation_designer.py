import os
import sys
import logging
import traceback
from typing import Dict, Any, List, Optional, Tuple
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from digitalagedu.core.llm import slide_kit
from digitalagedu.core.llm.schemas import Module, ProblemStatementSchema


logger = logging.getLogger(__name__)


class PresentationCodeSchema(BaseModel):
    """Schema for LLM-synthesized Python presentation generation script."""
    design_rationale: str = Field(
        ...,
        description="High-level pedagogical and visual layout rationale for the 5-slide deck."
    )
    python_code: str = Field(
        ...,
        description="Complete, valid Python script defining `def build_presentation(prs, slide_kit, telemetry, solution_code):` using python-pptx and slide_kit."
    )


def build_presentation_designer_prompt(
    module: Module,
    problem_formulation: ProblemStatementSchema,
    solution_code: str,
    telemetry: Dict[str, Any]
) -> str:
    """Builds the comprehensive instruction prompt for the LLM Presentation Designer."""
    title = problem_formulation.title or module.title
    domain_ctx = problem_formulation.domain_context or module.context
    problem_stmt = problem_formulation.problem_statement or ""
    objectives = problem_formulation.learning_objectives or []
    target_in = problem_formulation.target_input_shape or "[Batch, Channels, Height, Width]"
    target_out = problem_formulation.target_output_shape or "[Batch, NumClasses]"
    focus = problem_formulation.suggested_focus or "Deep Learning & Model Architecture"

    run_sum = telemetry.get("run_summary", {})
    acc = run_sum.get("overall_accuracy") or run_sum.get("accuracy", "90.57")
    auc = run_sum.get("auc_roc", "0.9709")
    samples = run_sum.get("total_images") or run_sum.get("total_samples", "3297")
    classes = telemetry.get("class_mapping", {})
    class_str = ", ".join(list(classes.values())[:6]) if classes else "Target Classes"

    contrastive = telemetry.get("contrastive_samples", {})
    top_succ = contrastive.get("top_success") or {}
    hard_fail = contrastive.get("hard_failure") or {}

    prompt = f"""You are an expert AI Presentation Designer and Applied Deep Learning Educator.
Your task is to write a clean, creative Python function `build_presentation(prs, slide_kit, telemetry, solution_code)` that creates a comprehensive, beautifully paced 16:9 widescreen PowerPoint lecture deck.

--- SLIDE COUNT & PACING AUTONOMY ---
You have full creative freedom to decide how many slides to generate based on the topic's depth (typically between 4 to 8 slides). Pacing matters: break complex concepts across multiple focused slides rather than overcrowding a single slide.

--- MODULE CONTEXT & TELEMETRY ---
Module: {module.title} (Week {module.week}) | Difficulty: {module.difficulty.title()}
Domain Application: {domain_ctx}
Core Problem Directive: {problem_stmt}
Focus Area: {focus}
Tensor Shape Contracts: Input `{target_in}` ➔ Output `{target_out}`
Learning Objectives:
{chr(10).join([f"- {obj}" for obj in objectives])}

Pipeline Telemetry:
- Total Dataset Samples: {samples}
- Phase 1 Baseline Accuracy: {acc}% | AUC-ROC: {auc}
- Dataset Classes: {class_str}
- Top Success Case: {top_succ.get('image_path', 'sample.jpg')} (True: {top_succ.get('ground_truth', 'Target')})
- Hard Failure Case: {hard_fail.get('image_path', 'failure.jpg')} (True: {hard_fail.get('ground_truth')}, Predicted: {hard_fail.get('predicted_class')})

Reference Solution Code:
```python
{solution_code}
```

--- AVAILABLE `slide_kit` UI TOOLKIT ---
You have access to `slide_kit` (and standard `python-pptx` `Inches`, `Pt`, `RGBColor`):
- `slide_kit.Theme`: `BG_DARK`, `BG_LIGHT`, `CARD_BG`, `CARD_BORDER`, `TEXT_PRIMARY`, `TEXT_MUTED`, `ACCENT_CYAN`, `ACCENT_INDIGO`, `ACCENT_EMERALD`, `ACCENT_CORAL`, `ACCENT_GOLD`, `CODE_BG`, `CODE_TEXT`.
- `slide_kit.create_slide(prs, bg_color=None)`: Creates a 16:9 slide (Width 13.333", Height 7.5").
- `slide_kit.add_header(slide, tag, title, subtitle=None, tag_color=None)`: Adds modern top header.
- `slide_kit.add_card(slide, x, y, w, h, title=None, body=None, bg_color=None, border_color=None, accent_color=None, title_size=12, body_size=12)`: Rounded container card.
- `slide_kit.add_code_box(slide, x, y, w, h, code, title=None, font_size=9.5)`: Dark IDE code block.
- `slide_kit.add_metric_card(slide, x, y, w, h, label, value, subtext=None, accent_color=None)`: KPI stat card.
- `slide_kit.add_badge_row(slide, x, y, badges=[('KEY', 'VAL'), ...], item_w=2.75, gap=0.24, h=0.9)`: Metadata chips row.
- `slide_kit.add_contrastive_cards(slide, x, y, w, h, success_dict, failure_dict)`: Side-by-side error analysis cards.
- `slide_kit.add_step_flow(slide, x, y, w, h, steps=['Step 1', 'Step 2', ...])`: Pipeline workflow diagram.
- `slide_kit.add_callout_banner(slide, x, y, w, h, text, title='KEY TAKEAWAY', accent_color=None)`: Takeaway banner.
- `slide_kit.add_table(slide, x, y, w, h, headers, rows)`: Formatted data grid.

--- RECOMMENDED TOPICS TO COVER (COMPOSE ACROSS SLIDES AS YOU SEE FIT) ---
- Hero Cover & Domain Context
- Real-World Challenge & Empirical Telemetry
- Computational Principles & Subsystem Contracts (Input ➔ Output)
- Reference Solution Architecture & Implementation Walkthrough
- Diagnostic Error Analysis & Contrastive Case Studies
- Pedagogical Takeaways & Production Considerations

--- CRITICAL CODE RULES ---
1. Define a top-level function: `def build_presentation(prs, slide_kit, telemetry, solution_code):`
2. Allowed imports at the top of your script:
   ```python
   import pptx
   from pptx.util import Inches, Pt
   from pptx.dml.color import RGBColor
   import slide_kit
   from slide_kit import Theme
   ```
3. Do NOT create a new `Presentation()` inside the function; use the passed `prs`.
4. Create each slide using `slide = slide_kit.create_slide(prs)`.
5. Ensure all coordinates fit within the 16:9 canvas (Width: 13.333 inches, Height: 7.5 inches).
6. Code MUST be valid, self-contained Python without markdown backticks in the python_code field.
"""
    return prompt



import re


class SafeTelemetryDict(dict):
    """
    A typo-tolerant, case-insensitive dictionary wrapper for telemetry.
    Guarantees that LLM code accessing keys like telemetry['Top Success Case'],
    telemetry['top_success'], telemetry['Phase 1 Baseline Accuracy'], etc.
    never raises KeyError or AttributeError.
    """
    def __init__(self, data=None):
        super().__init__()
        if data and isinstance(data, dict):
            for k, v in data.items():
                if isinstance(v, dict):
                    self[k] = SafeTelemetryDict(v)
                elif isinstance(v, list):
                    self[k] = [SafeTelemetryDict(item) if isinstance(item, dict) else item for item in v]
                else:
                    self[k] = v

            # Pre-populate common flat aliases derived from run_summary & contrastive_samples
            run_sum = data.get("run_summary", {})
            if isinstance(run_sum, dict):
                self["total_samples"] = run_sum.get("total_samples") or run_sum.get("total_images") or 0
                self["total_dataset_samples"] = self["total_samples"]
                self["accuracy"] = run_sum.get("overall_accuracy_percent") or run_sum.get("accuracy") or "90%"
                self["baseline_accuracy"] = self["accuracy"]
                self["phase_1_baseline_accuracy"] = self["accuracy"]
                self["auc_roc"] = run_sum.get("auc_roc") or "N/A"

            contrastive = data.get("contrastive_samples", {})
            if isinstance(contrastive, dict):
                for ck, cv in contrastive.items():
                    safe_cv = SafeTelemetryDict(cv) if isinstance(cv, dict) else cv
                    self[ck] = safe_cv
                    self[f"{ck}_case"] = safe_cv
                    self[f"{ck}_sample"] = safe_cv

    def __getitem__(self, key):
        if not isinstance(key, str):
            return super().get(key, SafeTelemetryDict())
        if key in self:
            return super().__getitem__(key)
        # Normalized key match (lowercase, remove spaces, underscores, hyphens)
        norm_key = re.sub(r"[\s_-]+", "", key.lower())
        for k, v in self.items():
            if isinstance(k, str) and re.sub(r"[\s_-]+", "", k.lower()) == norm_key:
                return v
        # Search inside nested dictionaries
        for sub in ["contrastive_samples", "run_summary"]:
            sub_dict = super().get(sub)
            if isinstance(sub_dict, dict):
                for sk, sv in sub_dict.items():
                    if isinstance(sk, str) and (sk.lower() == key.lower() or re.sub(r"[\s_-]+", "", sk.lower()) == norm_key):
                        return sv
        return SafeTelemetryDict()

    def get(self, key, default=None):
        try:
            val = self[key]
            if isinstance(val, SafeTelemetryDict) and len(val) == 0:
                return default
            return val
        except Exception:
            return default

    def __getattr__(self, name):
        return self[name]


def _build_default_presentation(
    prs: Presentation,
    module: Module,
    problem_formulation: ProblemStatementSchema,
    telemetry: Dict[str, Any],
    solution_code: str
) -> None:
    """Builds a clean, production-grade 4-slide fallback deck if LLM python script synthesis fails."""
    # Slide 1: Hero Cover
    s1 = slide_kit.create_slide(prs)
    slide_kit.add_header(s1, f"WEEK {module.week}", module.title, tag_color=slide_kit.Theme.ACCENT_CYAN)
    slide_kit.add_card(
        s1, 0.8, 2.0, 11.7, 4.8,
        title="Domain Context & Application",
        body=problem_formulation.domain_context or module.context or "Domain dataset analysis and algorithmic implementation."
    )

    # Slide 2: Telemetry & Challenge
    s2 = slide_kit.create_slide(prs)
    slide_kit.add_header(s2, "EMPIRICAL TELEMETRY", "Dataset Grounding & Challenge Directives", tag_color=slide_kit.Theme.ACCENT_INDIGO)
    run_sum = telemetry.get("run_summary", {}) if telemetry else {}
    acc = run_sum.get("overall_accuracy_percent") or run_sum.get("accuracy") or "90%"
    samples = run_sum.get("total_samples") or run_sum.get("total_images") or "3297"
    slide_kit.add_metric_card(s2, 0.8, 2.0, 3.6, 2.0, "Dataset Samples", str(samples), accent_color=slide_kit.Theme.ACCENT_CYAN)
    slide_kit.add_metric_card(s2, 4.8, 2.0, 3.6, 2.0, "Baseline Accuracy", f"{acc}%" if "%" not in str(acc) else str(acc), accent_color=slide_kit.Theme.ACCENT_EMERALD)
    slide_kit.add_metric_card(s2, 8.8, 2.0, 3.6, 2.0, "Academic Week", f"Week {module.week}", accent_color=slide_kit.Theme.ACCENT_GOLD)
    slide_kit.add_card(
        s2, 0.8, 4.4, 11.7, 2.4,
        title="Problem Statement",
        body=problem_formulation.problem_statement
    )

    # Slide 3: Architecture Walkthrough
    s3 = slide_kit.create_slide(prs)
    slide_kit.add_header(s3, "SOFTWARE ARCHITECTURE", "Reference Implementation Walkthrough", tag_color=slide_kit.Theme.ACCENT_GOLD)
    code_lines = [l for l in solution_code.strip().split("\n") if not l.startswith('"""')][:25]
    slide_kit.add_code_box(s3, 0.8, 2.0, 11.7, 4.8, "\n".join(code_lines), title="Core Module Pipeline")

    # Slide 4: Summary & Objectives
    s4 = slide_kit.create_slide(prs)
    slide_kit.add_header(s4, "KEY TAKEAWAYS", "Learning Outcomes & Production Practices", tag_color=slide_kit.Theme.ACCENT_EMERALD)
    objectives_text = "\n".join([f"• {obj}" for obj in (problem_formulation.learning_objectives or [])])
    slide_kit.add_card(s4, 0.8, 2.0, 11.7, 3.6, title="Target Learning Outcomes", body=objectives_text or "Applied computing literacy.")
    slide_kit.add_callout_banner(s4, 0.8, 6.0, 11.7, 0.9, "Hands-on student implementation completes all milestone subsystems.", title="SUMMARY")


def execute_presentation_code(
    python_code: str,
    prs: Presentation,
    telemetry: Dict[str, Any],
    solution_code: str
) -> Tuple[bool, Optional[str]]:
    """
    Executes the LLM-generated Python presentation script inside a controlled namespace with SafeTelemetryDict.
    """
    # Strip any potential markdown code fences
    cleaned_code = python_code.strip()
    if cleaned_code.startswith("```python"):
        cleaned_code = cleaned_code[len("```python"):].strip()
    elif cleaned_code.startswith("```"):
        cleaned_code = cleaned_code[3:].strip()
    if cleaned_code.endswith("```"):
        cleaned_code = cleaned_code[:-3].strip()

    # Register slide_kit and pptx in sys.modules for seamless top-level import support
    import pptx
    sys.modules["slide_kit"] = slide_kit
    sys.modules["pptx"] = pptx
    sys.modules["python_pptx"] = pptx

    safe_telemetry = SafeTelemetryDict(telemetry)

    exec_namespace = {
        "pptx": pptx,
        "Presentation": Presentation,
        "Inches": Inches,
        "Pt": Pt,
        "RGBColor": RGBColor,
        "slide_kit": slide_kit,
        "telemetry": safe_telemetry,
        "os": os,
        "sys": sys,
    }

    try:
        exec(cleaned_code, exec_namespace)
        build_func = exec_namespace.get("build_presentation")
        if not build_func or not callable(build_func):
            return False, "Generated code did not define a callable `build_presentation(prs, slide_kit, telemetry, solution_code)` function."

        build_func(prs, slide_kit, safe_telemetry, solution_code)
        return True, None
    except Exception as e:
        tb = traceback.format_exc()
        return False, f"{type(e).__name__}: {str(e)}\n\nTraceback:\n{tb}"


class PresentationDesigner:
    """
    Agentic Presentation Designer:
    Prompts vLLM to write custom python-pptx + slide_kit code, verifies execution in a sandbox,
    and self-heals syntax or layout errors with automatic retries and robust default deck fallback.
    """

    def __init__(self, client: Any = None, model_name: Optional[str] = None):
        self.client = client
        self.model_name = model_name or os.environ.get("LLM_MODEL", "Qwen/Qwen2.5-Coder-32B-Instruct-AWQ")

    def generate_presentation_deck(
        self,
        module: Module,
        problem_formulation: ProblemStatementSchema,
        solution_code: str,
        telemetry: Dict[str, Any],
        output_path: str
    ) -> str:
        """
        Generates and executes a custom 16:9 presentation deck via LLM code synthesis with self-healing
        and non-fatal default deck fallback.
        """
        prs = Presentation()
        prs.slide_width = Inches(slide_kit.SLIDE_WIDTH_INCHES)
        prs.slide_height = Inches(slide_kit.SLIDE_HEIGHT_INCHES)

        if not self.client:
            logger.info("No LLM client provided to PresentationDesigner. Generating default slide deck...")
            _build_default_presentation(prs, module, problem_formulation, telemetry, solution_code)
            os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
            prs.save(output_path)
            return output_path

        designer_prompt = build_presentation_designer_prompt(
            module=module,
            problem_formulation=problem_formulation,
            solution_code=solution_code,
            telemetry=telemetry
        )

        logger.info(f"Synthesizing custom 16:9 presentation code with {self.model_name}...")
        try:
            res: PresentationCodeSchema = self.client.chat.completions.create(
                model=self.model_name,
                response_model=PresentationCodeSchema,
                max_retries=2,
                max_tokens=6144,
                messages=[
                    {"role": "system", "content": "You are a master presentation designer and senior Python programmer."},
                    {"role": "user", "content": designer_prompt}
                ]
            )

            success, err = execute_presentation_code(res.python_code, prs, telemetry, solution_code)
            if not success:
                logger.warning(f"Presentation execution failed on initial attempt: {err}. Triggering self-healing retry...")
                retry_prompt = f"{designer_prompt}\n\n--- PREVIOUS EXECUTION ERROR ---\n{err}\n\nPlease fix the python_code to resolve the error."
                prs = Presentation()
                prs.slide_width = Inches(slide_kit.SLIDE_WIDTH_INCHES)
                prs.slide_height = Inches(slide_kit.SLIDE_HEIGHT_INCHES)

                retry_res: PresentationCodeSchema = self.client.chat.completions.create(
                    model=self.model_name,
                    response_model=PresentationCodeSchema,
                    max_retries=2,
                    max_tokens=6144,
                    messages=[
                        {"role": "system", "content": "You are a master presentation designer and senior Python programmer."},
                        {"role": "user", "content": retry_prompt}
                    ]
                )
                success, err = execute_presentation_code(retry_res.python_code, prs, telemetry, solution_code)
                if not success:
                    logger.warning(f"Failed to generate valid presentation code after retry: {err}. Falling back to default slide deck.")
                    prs = Presentation()
                    prs.slide_width = Inches(slide_kit.SLIDE_WIDTH_INCHES)
                    prs.slide_height = Inches(slide_kit.SLIDE_HEIGHT_INCHES)
                    _build_default_presentation(prs, module, problem_formulation, telemetry, solution_code)
        except Exception as e:
            logger.warning(f"Presentation synthesis hit exception: {e}. Falling back to default slide deck.")
            prs = Presentation()
            prs.slide_width = Inches(slide_kit.SLIDE_WIDTH_INCHES)
            prs.slide_height = Inches(slide_kit.SLIDE_HEIGHT_INCHES)
            _build_default_presentation(prs, module, problem_formulation, telemetry, solution_code)

        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        prs.save(output_path)
        logger.info(f"Successfully saved 16:9 presentation deck: {output_path} ({os.path.getsize(output_path)} bytes)")
        return output_path

