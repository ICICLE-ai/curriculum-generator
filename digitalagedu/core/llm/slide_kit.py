import os
import logging
from typing import Dict, Any, List, Optional, Tuple, Union
from pptx import Presentation
from pptx.util import Inches, Pt, Length
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

# Standard 16:9 Widescreen dimensions
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5


def _to_length(val: Union[int, float, Length]) -> Length:
    """
    Safely converts a coordinate or dimension value to a pptx Length.
    - If already a Length object or a large integer (> 1000, i.e. already in EMUs), returns as-is.
    - If a float or small integer (<= 1000), treats it as inches and wraps with Inches(val).
    - Prevents double Inches(Inches(x)) multiplication that causes trillion-EMU coordinate overflows in PowerPoint.
    """
    if isinstance(val, Length):
        return val
    if isinstance(val, (int, float)):
        if val > 1000:
            return Length(int(val))
        return Inches(float(val))
    try:
        f_val = float(val)
        if f_val > 1000:
            return Length(int(f_val))
        return Inches(f_val)
    except (ValueError, TypeError):
        return Inches(1.0)


class _ThemeMeta(type):
    """Metaclass providing graceful typo tolerance for Theme attribute access."""
    def __getattr__(cls, name: str) -> RGBColor:
        clean = name.upper().replace("_", "")
        # Try direct or substring match against known colors
        for key, val in cls.__dict__.items():
            if isinstance(val, RGBColor) and key.startswith("ACCENT_"):
                if key.replace("ACCENT_", "").replace("_", "") in clean or clean in key.replace("_", ""):
                    return val
        if "BG" in clean:
            return cls.BG_DARK
        if "TEXT" in clean or "WHITE" in clean:
            return cls.TEXT_PRIMARY
        if "MUTED" in clean or "GRAY" in clean or "GREY" in clean:
            return cls.TEXT_MUTED
        if "CODE" in clean:
            return cls.CODE_BG
        # Default safe high-visibility cyan accent
        return cls.ACCENT_CYAN


class Theme(metaclass=_ThemeMeta):
    """Pre-calibrated 16:9 modern presentation color palette with typo-tolerant attribute access."""
    BG_DARK = RGBColor(11, 15, 25)           # #0B0F19 Deep slate midnight
    BG_LIGHT = RGBColor(248, 250, 252)       # #F8FAFC Clean studio light
    CARD_BG = RGBColor(30, 41, 59)          # #1E293B Card container background
    CARD_BORDER = RGBColor(51, 65, 85)       # #334155 Card border stroke
    TEXT_PRIMARY = RGBColor(248, 250, 252)   # #F8FAFC High-contrast white
    TEXT_MUTED = RGBColor(148, 163, 184)     # #94A3B8 Secondary cool gray
    ACCENT_CYAN = RGBColor(56, 189, 248)     # #38BDF8 Electric Cyan
    ACCENT_INDIGO = RGBColor(99, 102, 241)   # #6366F1 Modern Indigo
    ACCENT_EMERALD = RGBColor(16, 185, 129)  # #10B981 Success Emerald
    ACCENT_CORAL = RGBColor(239, 68, 68)     # #EF4444 Diagnostic Coral / Alert
    ACCENT_GOLD = RGBColor(245, 158, 11)     # #F59E0B Highlight Amber
    CODE_BG = RGBColor(15, 23, 42)           # #0F172A Dark IDE Canvas
    CODE_TEXT = RGBColor(226, 232, 240)      # #E2E8F0 Monospaced code text


def create_slide(prs: Presentation, bg_color: Optional[RGBColor] = None, *args, **kwargs):
    """Creates a new blank 16:9 slide with a full-bleed solid background."""
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg_col = bg_color or kwargs.get("bg") or kwargs.get("color") or Theme.BG_DARK
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, Inches(SLIDE_WIDTH_INCHES), Inches(SLIDE_HEIGHT_INCHES)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_col
    bg.line.fill.background()
    return slide


def add_header(
    slide,
    tag: str = "CURRICULUM",
    title: str = "",
    subtitle: Optional[str] = None,
    tag_color: Optional[RGBColor] = None,
    *args,
    **kwargs
):
    """Adds a standardized top category pill, bold title, and optional subtitle."""
    tg = str(tag or kwargs.get("category") or kwargs.get("topic") or "CURRICULUM")
    ttl = str(title or kwargs.get("heading") or kwargs.get("title_text") or "")
    sub = subtitle or kwargs.get("sub") or kwargs.get("description")
    accent = tag_color or kwargs.get("accent") or Theme.ACCENT_CYAN
    
    # 1. Category pill
    pill_w = max(2.2, len(tg) * 0.11 + 0.5)
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        _to_length(0.8), _to_length(0.5), _to_length(pill_w), _to_length(0.38)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = Theme.CARD_BG
    pill.line.color.rgb = accent
    pill.line.width = Pt(1.2)

    tf_pill = pill.text_frame
    tf_pill.word_wrap = False
    tf_pill.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_pill = tf_pill.paragraphs[0]
    p_pill.text = tg.upper()
    p_pill.font.size = Pt(9.5)
    p_pill.font.bold = True
    p_pill.font.color.rgb = accent
    p_pill.alignment = PP_ALIGN.CENTER

    # 2. Main title & subtitle
    tbox = slide.shapes.add_textbox(
        _to_length(0.8), _to_length(0.95), _to_length(11.733), _to_length(0.9)
    )
    tf = tbox.text_frame
    tf.word_wrap = True
    p_title = tf.paragraphs[0]
    p_title.text = ttl
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = Theme.TEXT_PRIMARY

    if sub:
        p_sub = tf.add_paragraph()
        p_sub.text = str(sub)
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = Theme.TEXT_MUTED
        p_sub.space_before = Pt(3)


def add_card(
    slide,
    x: Union[int, float, Length],
    y: Union[int, float, Length],
    w: Union[int, float, Length],
    h: Union[int, float, Length],
    title: Optional[str] = None,
    body: Optional[str] = None,
    bg_color: Optional[RGBColor] = None,
    border_color: Optional[RGBColor] = None,
    accent_color: Optional[RGBColor] = None,
    title_size: int = 12,
    body_size: int = 12,
    *args,
    **kwargs
):
    """Adds a rounded rectangular card container with optional title and body paragraphs."""
    t = title or kwargs.get("heading") or kwargs.get("header")
    b = body or kwargs.get("text") or kwargs.get("content") or kwargs.get("description")
    bg = bg_color or kwargs.get("bg") or Theme.CARD_BG
    border = border_color or kwargs.get("border") or Theme.CARD_BORDER
    acc = accent_color or kwargs.get("accent") or Theme.ACCENT_CYAN

    lx, ly, lw, lh = _to_length(x), _to_length(y), _to_length(w), _to_length(h)

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        lx, ly, lw, lh
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = border
    card.line.width = Pt(1.5)

    if t or b:
        pad_x = _to_length(0.2)
        pad_y = _to_length(0.18)
        tbox = slide.shapes.add_textbox(
            lx + pad_x, ly + pad_y,
            max(_to_length(0.5), lw - (pad_x * 2)),
            max(_to_length(0.5), lh - (pad_y * 2))
        )
        tf = tbox.text_frame
        tf.word_wrap = True

        first = True
        if t:
            p_t = tf.paragraphs[0]
            p_t.text = str(t)
            p_t.font.size = Pt(title_size)
            p_t.font.bold = True
            p_t.font.color.rgb = acc
            first = False

        if b:
            p_b = tf.paragraphs[0] if first else tf.add_paragraph()
            p_b.text = str(b)
            p_b.font.size = Pt(body_size)
            p_b.font.color.rgb = Theme.TEXT_PRIMARY
            if not first:
                p_b.space_before = Pt(6)

    return card


def add_code_box(
    slide,
    x: Union[int, float, Length],
    y: Union[int, float, Length],
    w: Union[int, float, Length],
    h: Union[int, float, Length],
    code: Optional[str] = None,
    title: Optional[str] = None,
    font_size: float = 9.5,
    *args,
    **kwargs
):
    """Renders a syntax-styled dark IDE code block with monospaced font."""
    raw_code = code or kwargs.get("code_string") or kwargs.get("code_str") or kwargs.get("text") or ""
    t = title or kwargs.get("heading") or kwargs.get("header")

    lx, ly, lw, lh = _to_length(x), _to_length(y), _to_length(w), _to_length(h)

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        lx, ly, lw, lh
    )
    card.fill.solid()
    card.fill.fore_color.rgb = Theme.CODE_BG
    card.line.color.rgb = Theme.CARD_BORDER
    card.line.width = Pt(1.5)

    pad_x = _to_length(0.2)
    pad_y = _to_length(0.15)
    tbox = slide.shapes.add_textbox(
        lx + pad_x, ly + pad_y,
        max(_to_length(0.5), lw - (pad_x * 2)),
        max(_to_length(0.5), lh - (pad_y * 2))
    )
    tf = tbox.text_frame
    tf.word_wrap = True

    first = True
    if t:
        p_t = tf.paragraphs[0]
        p_t.text = str(t).upper()
        p_t.font.size = Pt(10)
        p_t.font.bold = True
        p_t.font.color.rgb = Theme.ACCENT_CYAN
        first = False

    lines = [l for l in str(raw_code).strip().split("\n") if not l.startswith('"""') and not l.startswith("'''")]
    display_code = "\n".join(lines[:24]) if lines else str(raw_code)

    p_c = tf.paragraphs[0] if first else tf.add_paragraph()
    p_c.text = display_code
    p_c.font.name = "Consolas"
    p_c.font.size = Pt(font_size)
    p_c.font.color.rgb = Theme.CODE_TEXT
    if not first:
        p_c.space_before = Pt(4)

    return card


def add_metric_card(
    slide,
    x: Union[int, float, Length],
    y: Union[int, float, Length],
    w: Union[int, float, Length],
    h: Union[int, float, Length],
    label: Optional[str] = None,
    value: Optional[Any] = None,
    subtext: Optional[str] = None,
    accent_color: Optional[RGBColor] = None,
    *args,
    **kwargs
):
    """Renders a high-impact KPI statistic card."""
    lbl = label or kwargs.get("title") or kwargs.get("name") or ""
    val = value if value is not None else (kwargs.get("val") or kwargs.get("metric") or "")
    sub = subtext or kwargs.get("subtitle") or kwargs.get("sub") or kwargs.get("description")
    acc = accent_color or kwargs.get("accent") or Theme.ACCENT_CYAN

    lx, ly, lw, lh = _to_length(x), _to_length(y), _to_length(w), _to_length(h)

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        lx, ly, lw, lh
    )
    card.fill.solid()
    card.fill.fore_color.rgb = Theme.CARD_BG
    card.line.color.rgb = Theme.CARD_BORDER
    card.line.width = Pt(1.5)

    pad_x = _to_length(0.15)
    pad_y = _to_length(0.12)
    tbox = slide.shapes.add_textbox(
        lx + pad_x, ly + pad_y,
        max(_to_length(0.5), lw - (pad_x * 2)),
        max(_to_length(0.5), lh - (pad_y * 2))
    )
    tf = tbox.text_frame
    tf.word_wrap = True

    # Label
    p_lbl = tf.paragraphs[0]
    p_lbl.text = str(lbl).upper()
    p_lbl.font.size = Pt(9)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = Theme.TEXT_MUTED

    # Value
    p_val = tf.add_paragraph()
    p_val.text = str(val)
    p_val.font.size = Pt(20)
    p_val.font.bold = True
    p_val.font.color.rgb = acc
    p_val.space_before = Pt(2)

    # Subtext
    if sub:
        p_sub = tf.add_paragraph()
        p_sub.text = str(sub)
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = Theme.TEXT_MUTED
        p_sub.space_before = Pt(2)

    return card


def add_badge_row(
    slide,
    x: Union[int, float, Length],
    y: Union[int, float, Length],
    badges: Optional[List[Tuple[str, str]]] = None,
    item_w: Union[int, float, Length] = 2.75,
    gap: Union[int, float, Length] = 0.24,
    h: Union[int, float, Length] = 0.9,
    *args,
    **kwargs
):
    """Renders an evenly spaced horizontal row of metadata chips."""
    bdgs = badges or kwargs.get("items") or kwargs.get("chips") or []
    lx, ly = _to_length(x), _to_length(y)
    liw, lgap, lh = _to_length(item_w), _to_length(gap), _to_length(h)

    for i, (b_title, b_val) in enumerate(bdgs):
        bx = lx + (i * (liw + lgap))
        add_metric_card(
            slide,
            x=bx, y=ly, w=liw, h=lh,
            label=str(b_title),
            value=str(b_val),
            accent_color=Theme.ACCENT_CYAN
        )


def add_contrastive_cards(
    slide,
    x: Union[int, float, Length],
    y: Union[int, float, Length],
    w: Union[int, float, Length],
    h: Union[int, float, Length],
    success_data: Optional[Dict[str, Any]] = None,
    failure_data: Optional[Dict[str, Any]] = None,
    *args,
    **kwargs
):
    """
    Renders side-by-side diagnostic case studies (Top Success vs Hard Failure).
    Accepts any keyword alias (success_data, success_dict, success, top_success, failure_data, failure_dict, failure, hard_failure).
    """
    succ = (
        success_data or 
        kwargs.get("success_dict") or 
        kwargs.get("success") or 
        kwargs.get("top_success") or 
        (args[0] if len(args) > 0 else {}) or 
        {}
    )
    fail = (
        failure_data or 
        kwargs.get("failure_dict") or 
        kwargs.get("failure") or 
        kwargs.get("hard_failure") or 
        (args[1] if len(args) > 1 else {}) or 
        {}
    )

    lx, ly, lw, lh = _to_length(x), _to_length(y), _to_length(w), _to_length(h)
    gap = _to_length(0.3)
    col_w = (lw - gap) / 2.0

    # 1. Success Card (Emerald)
    add_card(
        slide,
        x=lx, y=ly, w=col_w, h=lh,
        title="HIGH-CONFIDENCE MATCH (SUCCESS)",
        body=(
            f"Sample: {os.path.basename(succ.get('image_path', succ.get('Sample Path', 'sample.jpg'))) if isinstance(succ, dict) else 'sample.jpg'}\n"
            f"Ground Truth: {succ.get('ground_truth', succ.get('True Category', 'Target')) if isinstance(succ, dict) else 'Target'}\n"
            f"Prediction: {succ.get('predicted_class', succ.get('ground_truth', 'Target')) if isinstance(succ, dict) else 'Target'}\n"
            f"Probabilities: {succ.get('probabilities', 'High Confidence') if isinstance(succ, dict) else 'High Confidence'}"
        ),
        accent_color=Theme.ACCENT_EMERALD,
        title_size=11,
        body_size=11
    )

    # 2. Failure Card (Coral Red)
    add_card(
        slide,
        x=lx + col_w + gap, y=ly, w=col_w, h=lh,
        title="DIAGNOSTIC FAILURE MODE (EDGE CASE)",
        body=(
            f"Sample: {os.path.basename(fail.get('image_path', fail.get('Sample Path', 'failure.jpg'))) if isinstance(fail, dict) else 'failure.jpg'}\n"
            f"Ground Truth: {fail.get('ground_truth', fail.get('True Category', 'True Class')) if isinstance(fail, dict) else 'True Class'}\n"
            f"Predicted: {fail.get('predicted_class', 'Misclassified') if isinstance(fail, dict) else 'Misclassified'}\n"
            f"Probabilities: {fail.get('probabilities', 'Shifted Decision Boundary') if isinstance(fail, dict) else 'Shifted Boundary'}"
        ),
        accent_color=Theme.ACCENT_CORAL,
        title_size=11,
        body_size=11
    )


def add_step_flow(
    slide,
    x: Union[int, float, Length],
    y: Union[int, float, Length],
    w: Union[int, float, Length],
    h: Union[int, float, Length],
    steps: Optional[List[str]] = None,
    *args,
    **kwargs
):
    """Renders a connected sequence of pipeline stages."""
    stps = steps or kwargs.get("pipeline_steps") or kwargs.get("items") or []
    n = len(stps)
    if n == 0:
        return

    lx, ly, lw, lh = _to_length(x), _to_length(y), _to_length(w), _to_length(h)
    gap = _to_length(0.2)
    step_w = (lw - (gap * (n - 1))) / n

    for i, step_text in enumerate(stps):
        sx = lx + (i * (step_w + gap))
        add_card(
            slide,
            x=sx, y=ly, w=step_w, h=lh,
            title=f"STAGE {i+1}",
            body=str(step_text),
            accent_color=Theme.ACCENT_CYAN,
            title_size=10,
            body_size=11
        )


def add_callout_banner(
    slide,
    x: Union[int, float, Length],
    y: Union[int, float, Length],
    w: Union[int, float, Length],
    h: Union[int, float, Length],
    text: Optional[str] = None,
    title: str = "KEY PEDAGOGICAL TAKEAWAY",
    accent_color: Optional[RGBColor] = None,
    *args,
    **kwargs
):
    """Renders a full-width takeaway callout card with accent header."""
    txt = text or kwargs.get("body") or kwargs.get("content") or ""
    t = title or kwargs.get("heading") or "KEY PEDAGOGICAL TAKEAWAY"
    acc = accent_color or kwargs.get("accent") or Theme.ACCENT_GOLD
    return add_card(
        slide,
        x=x, y=y, w=w, h=h,
        title=str(t),
        body=str(txt),
        accent_color=acc,
        title_size=10,
        body_size=11.5
    )


def add_table(
    slide,
    x: Union[int, float, Length],
    y: Union[int, float, Length],
    w: Union[int, float, Length],
    h: Union[int, float, Length],
    headers: Optional[List[str]] = None,
    rows: Optional[List[List[str]]] = None,
    *args,
    **kwargs
):
    """Renders a styled data comparison table on the slide."""
    hdrs = headers or kwargs.get("columns") or []
    rws = rows or kwargs.get("data") or []
    num_rows = len(rws) + 1
    num_cols = len(hdrs)

    lx, ly, lw, lh = _to_length(x), _to_length(y), _to_length(w), _to_length(h)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols, lx, ly, lw, lh
    )
    table = table_shape.table

    for c_idx, head in enumerate(hdrs):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = Theme.CARD_BG
        p = cell.text_frame.paragraphs[0]
        p.text = str(head).upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = Theme.ACCENT_CYAN
        p.alignment = PP_ALIGN.CENTER

    for r_idx, row in enumerate(rws):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = Theme.CODE_BG
            p = cell.text_frame.paragraphs[0]
            p.text = str(val)
            p.font.size = Pt(10)
            p.font.color.rgb = Theme.TEXT_PRIMARY
            p.alignment = PP_ALIGN.CENTER

    return table_shape
