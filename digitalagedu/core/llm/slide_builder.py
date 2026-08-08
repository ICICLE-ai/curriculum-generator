import os
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from digitalagedu.core.llm.schemas import SlideDeckSchema

# Color Palette (Dark Professional / Modern EdTech)
BG_COLOR = RGBColor(15, 23, 42)        # Slate 900
PANEL_BG = RGBColor(30, 41, 59)        # Slate 800
TEXT_MAIN = RGBColor(248, 250, 252)    # Slate 50
TEXT_MUTED = RGBColor(148, 163, 184)  # Slate 400
ACCENT_BLUE = RGBColor(56, 189, 248)   # Sky 400
ACCENT_GREEN = RGBColor(52, 211, 153)  # Emerald 400
CODE_BG = RGBColor(15, 23, 42)         # Slate 900

def _apply_background(slide, color=BG_COLOR):
    """Sets a solid dark background color on a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def build_pptx_deck(slide_deck: SlideDeckSchema, output_path: str):
    """Programmatically builds a native 16:9 widescreen PowerPoint deck using python-pptx."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]

    # 1. Title Slide
    title_slide = prs.slides.add_slide(blank_slide_layout)
    _apply_background(title_slide, BG_COLOR)

    txBox = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = slide_deck.deck_title
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = TEXT_MAIN
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = "DigitalAgEdu Applied Deep Learning Suite"
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT_BLUE
    p2.font.name = "Segoe UI"
    p2.alignment = PP_ALIGN.LEFT

    # 2. Content Slides
    for slide_data in slide_deck.slides:
        slide = prs.slides.add_slide(blank_slide_layout)
        _apply_background(slide, BG_COLOR)

        # Slide Title Header
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_data.title
        p_title.font.bold = True
        p_title.font.size = Pt(28)
        p_title.font.color.rgb = TEXT_MAIN
        p_title.font.name = "Segoe UI"

        # Content Layout logic: 2-column if code_snippet exists, 1-column if text only
        has_code = bool(slide_data.code_snippet and slide_data.code_snippet.strip())
        text_width = Inches(5.8) if has_code else Inches(11.733)

        # Bullet Points Panel
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), text_width, Inches(5.2))
        tf_body = text_box.text_frame
        tf_body.word_wrap = True

        for i, pt in enumerate(slide_data.bullet_points):
            p = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
            p.text = f"• {pt}"
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_MAIN
            p.font.name = "Segoe UI"
            p.space_after = Pt(14)

        # Optional Code Snippet Panel
        if has_code:
            code_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2))
            tf_code = code_box.text_frame
            tf_code.word_wrap = True
            
            p_code_header = tf_code.paragraphs[0]
            p_code_header.text = "PyTorch Implementation:"
            p_code_header.font.bold = True
            p_code_header.font.size = Pt(14)
            p_code_header.font.color.rgb = ACCENT_GREEN
            p_code_header.font.name = "Consolas"
            p_code_header.space_after = Pt(8)

            p_code = tf_code.add_paragraph()
            p_code.text = slide_data.code_snippet.strip()
            p_code.font.size = Pt(12)
            p_code.font.color.rgb = ACCENT_BLUE
            p_code.font.name = "Consolas"

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
